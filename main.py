from fastapi import FastAPI, HTTPException
from fastapi.responses import JSONResponse, FileResponse
from pydantic import BaseModel, Field
from typing import Optional, List, Dict, Any, Literal
from pathlib import Path
from pptx import Presentation
import shutil
import uuid
import unicodedata
import re
import zipfile
import xml.etree.ElementTree as ET
import os

from v27.bootstrap import build_v27_runtime_safely
from v27.middleware import V27OperationIdMiddleware
from v27.routes import build_v27_router
from v27.security import resolve_safe_output_path


app = FastAPI(
    title="Estate Gover Presentation Generator",
    version="2.6.0"
)

BASE_DIR = Path(__file__).resolve().parent
TEMPLATES_DIR = BASE_DIR / "templates"
OUTPUTS_DIR = BASE_DIR / "outputs"
OUTPUTS_DIR.mkdir(parents=True, exist_ok=True)

# URL pública oficial do engine.
# Pode ser sobrescrita por variável de ambiente no Render, se necessário.
PUBLIC_BASE_URL = os.getenv(
    "PUBLIC_BASE_URL",
    "https://estate-gover-pptx-engine.onrender.com"
).rstrip("/")

# Governança fixa do backend.
PUBLICACAO_AUTOMATICA = False

# Runtime isolado V2.7. Drive permanece desabilitado por padrão e qualquer
# falha de configuração recua com segurança sem indisponibilizar V1/V2.6.
V27_RUNTIME = build_v27_runtime_safely()
V27_AUDIT_LOGGER = V27_RUNTIME.audit_logger
app.add_middleware(
    V27OperationIdMiddleware,
    audit_logger=V27_AUDIT_LOGGER,
)
app.include_router(build_v27_router(
    V27_AUDIT_LOGGER,
    V27_RUNTIME.drive_provider,
    V27_RUNTIME.renderer,
    V27_RUNTIME.serial_guard,
))


# ============================================================
# MODELOS DE DADOS
# ============================================================

class PresentationRequest(BaseModel):
    """
    V1 — mantida apenas para rollback seguro em /v1.
    Não deve ser usada como rota principal.
    """

    codigo_ativo: str
    nome_ativo: str
    localizacao: str
    tipo_ativo: Optional[str] = None
    area_aproximada: Optional[str] = None
    valor_referencia: Optional[str] = None
    texto_base_capa: str
    texto_base_estrategica: str
    observacoes_urbanisticas: Optional[str] = None
    preservar_multimidia: bool = True
    formato_saida: str = "pptx"


class MediaItem(BaseModel):
    """
    Item de mídia V2.

    Regra:
    - usar somente mídia validada do ativo atual;
    - nunca aproveitar mídia de outro ativo;
    - mídia só é aplicada em shapes marcados como EG_SLOT_* ou EG_ASSET_*.
    """

    slot: str
    tipo: Optional[str] = None
    path: Optional[str] = None
    codigo_ativo: Optional[str] = None
    legenda: Optional[str] = None


class PresentationRequestV2(BaseModel):
    """
    V2.6 — substitui placeholders, aplica política de mídia asset_only
    e governa produto, versão e status do artefato.
    """

    codigo_ativo: str
    nome_ativo: str
    localizacao: str
    tipo_ativo: Optional[str] = None
    area_aproximada: Optional[str] = None
    valor_referencia: Optional[str] = None

    descricao_capa: Optional[str] = None
    resumo_executivo: Optional[str] = None
    conheca_ativo: Optional[str] = None
    localizacao_texto: Optional[str] = None
    dimensao_status: Optional[str] = None
    potencial_urbanistico: Optional[str] = None
    diferenciais: Optional[str] = None
    escala_ativo: Optional[str] = None
    condicoes_negocio: Optional[str] = None
    observacoes_urbanisticas: Optional[str] = None

    # Campos legados da V1, mantidos por compatibilidade com a Action.
    texto_base_capa: Optional[str] = None
    texto_base_estrategica: Optional[str] = None

    # Campos V2.
    midias: List[MediaItem] = Field(default_factory=list)
    media_policy: str = "asset_only"
    preservar_multimidia: bool = True
    formato_saida: str = "pptx"

    # Governança V2.6.
    tipo_saida: Literal["CAPA", "ESTRATEGICA", "AMBOS"] = "CAPA"
    versao_saida: str = "V01"
    status_arquivo: Literal["RASCUNHO", "APROVADO"] = "RASCUNHO"


PLACEHOLDER_FIELDS = [
    "codigo_ativo",
    "nome_ativo",
    "localizacao",
    "tipo_ativo",
    "area_aproximada",
    "valor_referencia",
    "descricao_capa",
    "resumo_executivo",
    "conheca_ativo",
    "localizacao_texto",
    "dimensao_status",
    "potencial_urbanistico",
    "diferenciais",
    "escala_ativo",
    "condicoes_negocio",
    "observacoes_urbanisticas",
    "texto_base_capa",
    "texto_base_estrategica",
]

FIXED_PREFIXES = (
    "EG_FIXED_",
    "EG_FIXED_LOGO",
    "EG_FIXED_QR",
    "EG_FIXED_WHATSAPP",
    "EG_FIXED_CONTATO",
)

ASSET_SLOT_PREFIXES = (
    "EG_SLOT_",
    "EG_ASSET_",
)

WHATSAPP_URL_RE = re.compile(
    r"https?://(?:api\.whatsapp\.com/send\?phone=|wa\.me/)[^\s<>\"]+",
    flags=re.IGNORECASE,
)

WHATSAPP_LABEL = "Abrir WhatsApp"

IMAGE_CTA_PATTERNS = (
    "aperte na imagem",
    "clique na imagem",
    "descubra o que te espera",
)


# ============================================================
# FUNÇÕES AUXILIARES
# ============================================================

def normalize_text(value: str) -> str:
    normalized = unicodedata.normalize("NFD", value or "")
    without_accents = normalized.encode("ascii", "ignore").decode("utf-8")
    return without_accents.lower().strip()


def normalize_filename(value: str) -> str:
    return normalize_text(value).replace(" ", "-")


def governed_filename_part(value: str) -> str:
    """
    Normaliza texto para nomenclatura governada de arquivo.

    Mantém:
    - letras e números;
    - vírgula, ponto, underscore e hífen.

    Remove acentos e troca demais separadores por hífen.
    """

    value = unicodedata.normalize("NFD", value or "")
    value = value.encode("ascii", "ignore").decode("utf-8")
    value = value.upper().strip()
    value = re.sub(r"[^A-Z0-9,._-]+", "-", value)
    value = re.sub(r"-{2,}", "-", value)
    return value.strip("-_.") or "ATIVO"


def normalize_version(value: str) -> str:
    """
    Aceita V1, V01, 1, 01 etc. e normaliza para V01, V02...
    """

    text = (value or "V01").strip().upper()
    match = re.fullmatch(r"V?(\d+)", text)

    if not match:
        raise HTTPException(
            status_code=400,
            detail="versao_saida deve seguir o padrão V01, V02, V03..."
        )

    number = int(match.group(1))

    if number < 1:
        raise HTTPException(
            status_code=400,
            detail="versao_saida deve ser igual ou superior a V01."
        )

    return f"V{number:02d}"


def build_asset_base_name(payload: PresentationRequestV2) -> str:
    """
    Evita duplicar o código EG quando nome_ativo já começa pelo código.
    """

    codigo = governed_filename_part(payload.codigo_ativo)
    nome = governed_filename_part(payload.nome_ativo)

    if nome == codigo or nome.startswith(f"{codigo}-") or nome.startswith(f"{codigo}_"):
        return nome

    return f"{codigo}_{nome}"


def build_output_filename(
    payload: PresentationRequestV2,
    kind: Literal["CAPA", "ESTRATEGICA"]
) -> str:
    version = normalize_version(payload.versao_saida)
    base_name = build_asset_base_name(payload)

    return (
        f"{base_name}_{kind}_{version}_{payload.status_arquivo}.pptx"
    )


def build_text_filename(payload: PresentationRequestV2) -> str:
    version = normalize_version(payload.versao_saida)
    base_name = build_asset_base_name(payload)

    return (
        f"{base_name}_{payload.tipo_saida}_TEXTO-BASE_"
        f"{version}_{payload.status_arquivo}.txt"
    )


def output_relative_url(path: Path) -> str:
    return f"/outputs/{path.name}"


def output_absolute_url(path: Path) -> str:
    return f"{PUBLIC_BASE_URL}{output_relative_url(path)}"


def output_descriptor(path: Path) -> Dict[str, str]:
    return {
        "arquivo": path.name,
        "url_relativa": output_relative_url(path),
        "url": output_absolute_url(path),
    }


def find_template(kind: str) -> Path:
    """
    Localiza os PPTX oficiais dentro da pasta templates.

    Regra V2:
    - usar somente nomes Estate Gover;
    - não aceitar "estate-governador" como fallback silencioso.
    """

    if kind == "capa":
        candidates = [
            "estate-gover-capa-modelo-final.pptx",
        ]
        keywords = ["estate-gover", "capa"]

    elif kind == "estrategica":
        candidates = [
            "estate-gover-estrategica-modelo-final.pptx",
            "estate-gover-estratégica-modelo-final.pptx",
        ]
        keywords = ["estate-gover", "estrateg"]

    else:
        raise ValueError("Tipo de modelo inválido.")

    for filename in candidates:
        path = TEMPLATES_DIR / filename
        if path.exists():
            return path

    for path in TEMPLATES_DIR.glob("*.pptx"):
        normalized = normalize_filename(path.name)

        if "governador" in normalized:
            continue

        if all(keyword in normalized for keyword in keywords):
            return path

    raise HTTPException(
        status_code=500,
        detail=(
            f"Modelo oficial {kind} não encontrado na pasta templates "
            "com nomenclatura Estate Gover."
        )
    )


def safe_value(
    value: Any,
    fallback: str = "[pendente de confirmação]"
) -> str:
    if value is None:
        return fallback

    text = str(value).strip()

    if not text:
        return fallback

    return text


def payload_to_dict(payload: BaseModel) -> Dict[str, Any]:
    """
    Compatível com Pydantic v1 e v2.
    """

    if hasattr(payload, "model_dump"):
        return payload.model_dump()

    return payload.dict()


def build_placeholder_mapping(
    payload: PresentationRequestV2
) -> Dict[str, str]:
    data = payload_to_dict(payload)

    return {
        f"{{{{{field}}}}}": safe_value(data.get(field))
        for field in PLACEHOLDER_FIELDS
    }


def replace_placeholders_in_presentation(
    prs: Presentation,
    mapping: Dict[str, str]
) -> int:
    """
    Substitui placeholders dentro de caixas de texto existentes.

    Não cria slides, não altera master, não muda layout, não remove QR,
    WhatsApp, hyperlinks, logos, rodapés ou objetos interativos.
    """

    replacements = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue

            if shape.text_frame is None:
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    original = run.text

                    for placeholder, value in mapping.items():
                        if placeholder in run.text:
                            run.text = run.text.replace(
                                placeholder,
                                value
                            )

                    if run.text != original:
                        replacements += 1

    return replacements


def is_fixed_shape(shape) -> bool:
    name = getattr(shape, "name", "") or ""
    upper_name = name.upper()

    return any(
        upper_name.startswith(prefix.upper())
        for prefix in FIXED_PREFIXES
    )


def get_asset_slot_name(shape) -> Optional[str]:
    name = getattr(shape, "name", "") or ""

    for prefix in ASSET_SLOT_PREFIXES:
        if name.upper().startswith(prefix.upper()):
            return name[len(prefix):]

    return None


def clear_text_shape(shape, message: str = "") -> bool:
    if not hasattr(shape, "text_frame"):
        return False

    if shape.text_frame is None:
        return False

    shape.text_frame.clear()

    if message:
        shape.text = message

    return True


def apply_media_policy(
    prs: Presentation,
    payload: PresentationRequestV2,
    warnings: List[str]
) -> Dict[str, int]:
    """
    Política de mídia segura para integridade do PowerPoint.

    Regras:
    - EG_FIXED_* deve ser preservado;
    - EG_SLOT_* ou EG_ASSET_* representa mídia do ativo;
    - nenhum shape é removido;
    - sem mídia válida, o slot original do modelo permanece;
    - mídia válida é aplicada por cima do espaço existente;
    - nenhuma mídia de outro ativo pode ser aplicada.
    """

    stats = {
        "fixed_preserved": 0,
        "asset_slots_found": 0,
        "asset_slots_replaced": 0,
        "asset_slots_neutralized": 0,
        "untagged_media_untouched": 0,
        "cover_slots_preserved_as_template": 0,
        "asset_slots_preserved_as_template": 0,
        "non_destructive_mode": 1,
    }

    midias_por_slot: Dict[str, MediaItem] = {}

    for item in payload.midias:
        if (
            item.codigo_ativo
            and item.codigo_ativo != payload.codigo_ativo
        ):
            warnings.append(
                f"Mídia ignorada no slot {item.slot}: "
                "código do ativo diferente."
            )
            continue

        midias_por_slot[item.slot] = item

    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape in list(slide.shapes):
            if is_fixed_shape(shape):
                stats["fixed_preserved"] += 1
                continue

            slot = get_asset_slot_name(shape)

            if not slot:
                continue

            stats["asset_slots_found"] += 1
            media = midias_por_slot.get(slot)

            if media and media.path:
                media_path = Path(media.path)

                if not media_path.exists():
                    warnings.append(
                        f"Mídia do slot {slot} não encontrada "
                        f"no caminho informado: {media.path}"
                    )

                    stats["asset_slots_preserved_as_template"] += 1

                    if slide_index == 1:
                        stats[
                            "cover_slots_preserved_as_template"
                        ] += 1

                    continue

                slide.shapes.add_picture(
                    str(media_path),
                    shape.left,
                    shape.top,
                    width=shape.width,
                    height=shape.height,
                )

                stats["asset_slots_replaced"] += 1

            else:
                stats["asset_slots_preserved_as_template"] += 1

                if slide_index == 1:
                    stats[
                        "cover_slots_preserved_as_template"
                    ] += 1

    if stats["asset_slots_found"] == 0:
        warnings.append(
            "Nenhum slot de mídia EG_SLOT_* ou EG_ASSET_* foi encontrado. "
            "A V2.6 não removeu mídia visual sem marcação para evitar quebrar "
            "QR, WhatsApp, hyperlinks ou relações internas."
        )

    return stats


def sanitize_visible_whatsapp_links(
    prs: Presentation
) -> int:
    """
    Substitui URLs visíveis de WhatsApp por texto comercial limpo,
    preservando o hyperlink sempre que possível.
    """

    replacements = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if (
                not hasattr(shape, "text_frame")
                or shape.text_frame is None
            ):
                continue

            for paragraph in shape.text_frame.paragraphs:
                runs = list(paragraph.runs)

                if runs:
                    paragraph_text = "".join(
                        (run.text or "")
                        for run in runs
                    )

                    match = WHATSAPP_URL_RE.search(
                        paragraph_text
                    )

                    if not match:
                        continue

                    url = match.group(0)

                    runs[0].text = WHATSAPP_LABEL

                    for run in runs[1:]:
                        run.text = ""

                    try:
                        runs[0].hyperlink.address = url
                    except Exception:
                        pass

                    replacements += 1
                    continue

                current_text = getattr(
                    shape,
                    "text",
                    ""
                ) or ""

                match = WHATSAPP_URL_RE.search(
                    current_text
                )

                if match:
                    url = match.group(0)
                    shape.text = WHATSAPP_LABEL

                    try:
                        first_paragraph = (
                            shape.text_frame.paragraphs[0]
                        )
                        first_run = first_paragraph.runs[0]
                        first_run.hyperlink.address = url
                    except Exception:
                        pass

                    replacements += 1

    return replacements


def neutralize_image_ctas_without_media(
    prs: Presentation,
    media_stats: Dict[str, int]
) -> int:
    """
    Remove chamadas como "Aperte na imagem" quando nenhum slot
    de mídia do ativo foi preenchido.
    """

    if media_stats.get("asset_slots_replaced", 0) > 0:
        return 0

    replacements = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if (
                not hasattr(shape, "text_frame")
                or shape.text_frame is None
            ):
                continue

            original_text = getattr(
                shape,
                "text",
                ""
            ) or ""

            normalized = normalize_text(
                original_text
            )

            if not normalized:
                continue

            if any(
                pattern in normalized
                for pattern in IMAGE_CTA_PATTERNS
            ):
                clear_text_shape(shape, "")
                replacements += 1

    return replacements


def update_core_metadata(
    prs: Presentation,
    payload: PresentationRequestV2,
    presentation_kind: str
) -> bool:
    """
    Atualiza metadados internos para evitar herança de dados do ativo-base.
    """

    try:
        props = prs.core_properties

        kind_label = (
            "CAPA"
            if presentation_kind == "capa"
            else "ESTRATÉGICA"
        )

        props.title = (
            f"Estate Gover {kind_label} - "
            f"{payload.codigo_ativo}"
        )

        props.subject = (
            f"{payload.nome_ativo} | "
            f"{payload.localizacao}"
        )

        props.author = "Estate Gover"
        props.last_modified_by = (
            "Estate Gover PPTX Engine"
        )
        props.category = (
            "Estate Gover Presentation"
        )
        props.keywords = (
            f"Estate Gover, "
            f"{payload.codigo_ativo}, "
            f"{kind_label}"
        )
        props.comments = (
            "Arquivo gerado pelo Estate Gover PPTX Engine "
            "a partir dos modelos oficiais. "
            "CAPA e ESTRATÉGICA permanecem produtos separados."
        )

        return True

    except Exception:
        return False


def sanitize_hidden_template_references(
    path: Path,
    payload: PresentationRequestV2,
    presentation_kind: str,
    warnings: List[str],
) -> int:
    """
    Remove referências técnicas ocultas herdadas do modelo/base,
    sem alterar layout, mídia, relações, QR, WhatsApp ou conteúdo visível.
    """

    sanitized = 0

    try:
        tmp_path = path.with_suffix(
            ".hidden-clean.tmp.pptx"
        )

        kind_label = (
            "CAPA"
            if presentation_kind == "capa"
            else "ESTRATÉGICA"
        )

        replacement_value = (
            f"Estate Gover {kind_label} — "
            "slot preservado do modelo oficial "
            f"para {payload.codigo_ativo}"
        )

        attr_pattern = re.compile(
            r'\b(?P<attr>descr|title)="[^"]*'
            r'(?:eg0013|canela|governador|/mnt/data/)'
            r'[^"]*"',
            flags=re.IGNORECASE,
        )

        with zipfile.ZipFile(
            path,
            "r"
        ) as zin, zipfile.ZipFile(
            tmp_path,
            "w",
            zipfile.ZIP_DEFLATED
        ) as zout:

            for name in zin.namelist():
                data = zin.read(name)

                if name.endswith(".xml"):
                    try:
                        xml_text = data.decode(
                            "utf-8"
                        )
                    except UnicodeDecodeError:
                        zout.writestr(
                            name,
                            data
                        )
                        continue

                    def repl(match):
                        nonlocal sanitized
                        sanitized += 1

                        return (
                            f'{match.group("attr")}'
                            f'="{replacement_value}"'
                        )

                    xml_text = attr_pattern.sub(
                        repl,
                        xml_text
                    )

                    data = xml_text.encode(
                        "utf-8"
                    )

                zout.writestr(
                    name,
                    data
                )

        tmp_path.replace(path)

        return sanitized

    except Exception as exc:
        warnings.append(
            "Falha ao limpar referências ocultas "
            f"herdadas em {path.name}: {exc}"
        )

        return sanitized


def prune_unused_slide_relationships(
    path: Path,
    warnings: List[str]
) -> int:
    """
    Remove relações órfãs de slides depois das alterações.
    """

    pruned = 0

    try:
        tmp_path = path.with_suffix(
            ".tmp.pptx"
        )

        rel_ns = (
            "http://schemas.openxmlformats.org/"
            "package/2006/relationships"
        )

        r_ns = (
            "http://schemas.openxmlformats.org/"
            "officeDocument/2006/relationships"
        )

        ET.register_namespace(
            "",
            rel_ns
        )

        with zipfile.ZipFile(
            path,
            "r"
        ) as zin, zipfile.ZipFile(
            tmp_path,
            "w",
            zipfile.ZIP_DEFLATED
        ) as zout:

            names = zin.namelist()

            for name in names:
                data = zin.read(name)

                if (
                    name.startswith(
                        "ppt/slides/_rels/slide"
                    )
                    and name.endswith(
                        ".xml.rels"
                    )
                ):
                    slide_name = (
                        name.replace(
                            "ppt/slides/_rels/",
                            "ppt/slides/"
                        )
                        .replace(
                            ".rels",
                            ""
                        )
                    )

                    if slide_name in names:
                        slide_root = ET.fromstring(
                            zin.read(slide_name)
                        )

                        used_ids = set()

                        for elem in slide_root.iter():
                            for (
                                attr_name,
                                attr_value
                            ) in elem.attrib.items():

                                if (
                                    attr_name.startswith(
                                        "{" + r_ns + "}"
                                    )
                                    and attr_value
                                ):
                                    used_ids.add(
                                        attr_value
                                    )

                        rel_root = ET.fromstring(
                            data
                        )

                        removed_any = False

                        for rel in list(rel_root):
                            rid = rel.attrib.get(
                                "Id"
                            )

                            rel_type = rel.attrib.get(
                                "Type",
                                ""
                            )

                            is_prunable = (
                                rel_type.endswith(
                                    "/image"
                                )
                                or rel_type.endswith(
                                    "/video"
                                )
                                or rel_type.endswith(
                                    "/media"
                                )
                                or rel_type.endswith(
                                    "/hyperlink"
                                )
                            )

                            if (
                                is_prunable
                                and rid
                                and rid not in used_ids
                            ):
                                rel_root.remove(rel)
                                pruned += 1
                                removed_any = True

                        if removed_any:
                            data = ET.tostring(
                                rel_root,
                                encoding="utf-8",
                                xml_declaration=True,
                            )

                zout.writestr(
                    name,
                    data
                )

        tmp_path.replace(path)

        return pruned

    except Exception as exc:
        warnings.append(
            "Falha ao limpar relações órfãs "
            f"do PPTX {path.name}: {exc}"
        )

        return pruned


def validate_pptx_package(
    path: Path,
    warnings: List[str]
) -> bool:
    """
    Validação técnica básica do pacote PPTX.
    """

    try:
        with zipfile.ZipFile(path) as zf:
            bad_file = zf.testzip()

        if bad_file:
            warnings.append(
                "PPTX com arquivo interno "
                f"possivelmente corrompido: {bad_file}"
            )

            return False

        return True

    except Exception as exc:
        warnings.append(
            "Falha ao validar pacote PPTX "
            f"{path.name}: {exc}"
        )

        return False


def gerar_texto_base_v1(
    payload: PresentationRequest
) -> str:
    return f"""ESTATE GOVER - TEXTO BASE

Código: {payload.codigo_ativo}
Ativo: {payload.nome_ativo}
Localização: {payload.localizacao}
Tipo: {payload.tipo_ativo or ""}
Área aproximada: {payload.area_aproximada or ""}
Valor de referência: {payload.valor_referencia or ""}

=== CAPA ===
{payload.texto_base_capa}

=== ESTRATÉGICA ===
{payload.texto_base_estrategica}

=== OBSERVAÇÕES URBANÍSTICAS ===
{payload.observacoes_urbanisticas or ""}
"""


def gerar_texto_base_v2(
    payload: PresentationRequestV2
) -> str:
    """
    O texto-base respeita o mesmo tipo_saida solicitado.
    Não mistura CAPA e ESTRATÉGICA quando apenas um produto foi pedido.
    """

    version = normalize_version(
        payload.versao_saida
    )

    header = f"""ESTATE GOVER — TEXTO BASE V2.6

Código: {payload.codigo_ativo}
Ativo: {payload.nome_ativo}
Localização: {payload.localizacao}
Tipo de saída: {payload.tipo_saida}
Versão do artefato: {version}
Status do arquivo: {payload.status_arquivo}
Tipo: {payload.tipo_ativo or "[pendente de confirmação]"}
Área aproximada: {payload.area_aproximada or "[pendente de confirmação]"}
Valor de referência: {payload.valor_referencia or "[pendente de confirmação]"}
"""

    blocks: List[str] = []

    if payload.tipo_saida in (
        "CAPA",
        "AMBOS"
    ):
        blocks.append(
            f"""
=== CAPA ===

{payload.descricao_capa
 or payload.texto_base_capa
 or "[pendente de confirmação]"}
"""
        )

    if payload.tipo_saida in (
        "ESTRATEGICA",
        "AMBOS"
    ):
        blocks.append(
            f"""
=== ESTRATÉGICA ===

Resumo executivo:
{payload.resumo_executivo
 or payload.texto_base_estrategica
 or "[pendente de confirmação]"}

Conheça o ativo:
{payload.conheca_ativo
 or "[pendente de confirmação]"}

Localização:
{payload.localizacao_texto
 or "[pendente de confirmação]"}

Dimensão e status:
{payload.dimensao_status
 or "[pendente de confirmação]"}

Potencial construtivo e urbanístico:
{payload.potencial_urbanistico
 or "[necessário DM / consulta urbanística]"}

Diferenciais:
{payload.diferenciais
 or "[pendente de confirmação]"}

Escala do ativo:
{payload.escala_ativo
 or "[pendente de confirmação]"}

Condições de negócio:
{payload.condicoes_negocio
 or "[pendente de confirmação]"}

Observações urbanísticas:
{payload.observacoes_urbanisticas
 or "[necessário DM / consulta urbanística]"}
"""
        )

    return header + "\n".join(blocks)


def process_pptx_v2(
    template_path: Path,
    output_path: Path,
    payload: PresentationRequestV2,
    warnings: List[str],
    presentation_kind: str,
) -> Dict[str, Any]:
    """
    Duplica matriz oficial, substitui placeholders, aplica política de mídia
    e corrige metadados/textos auxiliares sem alterar identidade visual.
    """

    shutil.copyfile(
        template_path,
        output_path
    )

    prs = Presentation(
        str(output_path)
    )

    metadata_updated = update_core_metadata(
        prs,
        payload,
        presentation_kind
    )

    mapping = build_placeholder_mapping(
        payload
    )

    replacements = (
        replace_placeholders_in_presentation(
            prs,
            mapping
        )
    )

    media_stats = apply_media_policy(
        prs,
        payload,
        warnings
    )

    whatsapp_links_sanitized = (
        sanitize_visible_whatsapp_links(
            prs
        )
    )

    image_ctas_neutralized = (
        neutralize_image_ctas_without_media(
            prs,
            media_stats
        )
    )

    prs.save(
        str(output_path)
    )

    unused_relationships_pruned = (
        prune_unused_slide_relationships(
            output_path,
            warnings
        )
    )

    hidden_template_refs_sanitized = (
        sanitize_hidden_template_references(
            output_path,
            payload,
            presentation_kind,
            warnings,
        )
    )

    pptx_package_validated = (
        validate_pptx_package(
            output_path,
            warnings
        )
    )

    if replacements == 0:
        warnings.append(
            "Nenhum placeholder textual foi substituído "
            f"em {output_path.name}. "
            "Confirme se o PPTX possui placeholders oficiais."
        )

    if not metadata_updated:
        warnings.append(
            "Metadados internos não puderam ser "
            f"atualizados em {output_path.name}."
        )

    return {
        "template": template_path.name,
        "output": output_path.name,
        "url": output_absolute_url(output_path),
        "text_replacements": replacements,
        "media_stats": media_stats,
        "metadata_updated": metadata_updated,
        "whatsapp_links_sanitized": (
            whatsapp_links_sanitized
        ),
        "image_ctas_neutralized": (
            image_ctas_neutralized
        ),
        "unused_relationships_pruned": (
            unused_relationships_pruned
        ),
        "hidden_template_refs_sanitized": (
            hidden_template_refs_sanitized
        ),
        "pptx_package_validated": (
            pptx_package_validated
        ),
    }


def reserve_output_paths(
    paths: List[Path]
) -> None:
    """
    Reserva todos os destinos sem permitir sobrescrita silenciosa.

    Se algum arquivo já existir:
    - retorna HTTP 409;
    - nenhum arquivo existente é removido;
    - o cliente deve solicitar nova versão.
    """

    existing = [
        path.name
        for path in paths
        if path.exists()
    ]

    if existing:
        raise HTTPException(
            status_code=409,
            detail={
                "mensagem": (
                    "Arquivo de saída já existe; "
                    "sobrescrita não autorizada."
                ),
                "arquivos_existentes": existing,
                "acao_requerida": (
                    "Solicite nova versao_saida, "
                    "por exemplo V02."
                ),
            },
        )

    reserved: List[Path] = []

    try:
        for path in paths:
            path.touch(
                exist_ok=False
            )
            reserved.append(path)

    except FileExistsError:
        for path in reserved:
            path.unlink(
                missing_ok=True
            )

        raise HTTPException(
            status_code=409,
            detail=(
                "Arquivo de saída criado concorrentemente; "
                "sobrescrita não autorizada."
            ),
        )

    except Exception:
        for path in reserved:
            path.unlink(
                missing_ok=True
            )

        raise


# ============================================================
# ROTAS
# ============================================================

@app.get("/")
def healthcheck():
    return {
        "status": "ok",
        "service": (
            "Estate Gover Presentation Generator"
        ),
        "version": "2.6.0",
        "routes": [
            "/gerar-apresentacao-estate-gover",
            "/v1/gerar-apresentacao-estate-gover",
            "/v2/gerar-apresentacao-estate-gover",
            "/v2.7/gerar-apresentacao-estate-gover",
        ],
        "default_route": (
            "/gerar-apresentacao-estate-gover "
            "uses V2.6"
        ),
        "governanca": {
            "tipo_saida_default": "CAPA",
            "versao_saida_default": "V01",
            "status_arquivo_default": "RASCUNHO",
            "publicacao_automatica": False,
            "sobrescrita_silenciosa": False,
            "url_absoluta": True,
        },
        "v2_7": {
            "estado": (
                "FASE_2_HOMOLOG_CONFIGURADA"
                if V27_RUNTIME.drive_provider is not None
                else "FASE_1_IMPLEMENTADA"
            ),
            "geracao_disponivel": V27_RUNTIME.renderer is not None,
            "drive_disponivel": V27_RUNTIME.drive_provider is not None,
            "publicacao_automatica": False,
        },
    }


@app.get(
    "/outputs/{filename}",
    name="baixar_arquivo"
)
def baixar_arquivo(
    filename: str
):
    file_path = resolve_safe_output_path(
        OUTPUTS_DIR,
        filename,
    )

    if not file_path.exists():
        raise HTTPException(
            status_code=404,
            detail="Arquivo não encontrado."
        )

    return FileResponse(
        file_path
    )


@app.post(
    "/v1/gerar-apresentacao-estate-gover"
)
def gerar_apresentacao_estate_gover_v1(
    payload: PresentationRequest
):
    """
    V1 preservada para rollback.

    Regra:
    - nenhuma alteração funcional na V1;
    - mantém UUID e comportamento legado;
    - não é a rota principal da V2.6.
    """

    if payload.formato_saida != "pptx":
        raise HTTPException(
            status_code=400,
            detail=(
                "Apenas formato pptx é permitido."
            )
        )

    capa_template = find_template(
        "capa"
    )

    estrategica_template = find_template(
        "estrategica"
    )

    job_id = (
        f"{payload.codigo_ativo}_"
        f"{uuid.uuid4().hex[:8]}"
    )

    capa_out = (
        OUTPUTS_DIR
        / f"{job_id}_CAPA.pptx"
    )

    estrategica_out = (
        OUTPUTS_DIR
        / f"{job_id}_ESTRATEGICA.pptx"
    )

    shutil.copyfile(
        capa_template,
        capa_out
    )

    shutil.copyfile(
        estrategica_template,
        estrategica_out
    )

    txt_out = (
        OUTPUTS_DIR
        / f"{job_id}_texto_base.txt"
    )

    txt_out.write_text(
        gerar_texto_base_v1(payload),
        encoding="utf-8"
    )

    return JSONResponse({
        "status": "ok",
        "versao": "v1",
        "mensagem": (
            "Arquivos gerados a partir "
            "dos modelos oficiais."
        ),
        "capa_pptx_url": (
            f"/outputs/{capa_out.name}"
        ),
        "estrategica_pptx_url": (
            f"/outputs/{estrategica_out.name}"
        ),
        "texto_base_url": (
            f"/outputs/{txt_out.name}"
        ),
    })


@app.post(
    "/gerar-apresentacao-estate-gover"
)
@app.post(
    "/v2/gerar-apresentacao-estate-gover"
)
def gerar_apresentacao_estate_gover_v2(
    payload: PresentationRequestV2
):
    """
    V2.6 — rota principal.

    Regras:
    - gera CAPA, ESTRATÉGICA ou AMBOS;
    - default = CAPA;
    - AMBOS gera CAPA antes de ESTRATÉGICA;
    - versão do artefato é independente da arquitetura V03;
    - default da versão = V01;
    - status default = RASCUNHO;
    - sem UUID/hash nos arquivos V2.6;
    - sem sobrescrita silenciosa;
    - URL absoluta no retorno;
    - publicação automática é sempre falsa;
    - templates e V1 permanecem preservados.
    """

    if payload.formato_saida != "pptx":
        raise HTTPException(
            status_code=400,
            detail=(
                "Apenas formato pptx é permitido."
            )
        )

    # Valida/normaliza versão antes de reservar arquivos.
    version = normalize_version(
        payload.versao_saida
    )

    # Reatribui somente em memória para refletir retorno normalizado.
    payload.versao_saida = version

    warnings: List[str] = []

    selected_kinds = {
        "CAPA": ["capa"],
        "ESTRATEGICA": ["estrategica"],
        "AMBOS": ["capa", "estrategica"],
    }[payload.tipo_saida]

    output_paths: Dict[str, Path] = {}

    for kind in selected_kinds:
        kind_upper = (
            "CAPA"
            if kind == "capa"
            else "ESTRATEGICA"
        )

        output_paths[kind] = (
            OUTPUTS_DIR
            / build_output_filename(
                payload,
                kind_upper
            )
        )

    txt_out = (
        OUTPUTS_DIR
        / build_text_filename(payload)
    )

    reserved_paths = [
        *output_paths.values(),
        txt_out,
    ]

    reserve_output_paths(
        reserved_paths
    )

    arquivos: Dict[str, Any] = {}
    resultados: Dict[str, Any] = {}

    try:
        for kind in selected_kinds:
            template = find_template(
                kind
            )

            result = process_pptx_v2(
                template,
                output_paths[kind],
                payload,
                warnings,
                presentation_kind=kind,
            )

            resultados[kind] = result

            arquivos[
                f"{kind}_pptx"
            ] = output_descriptor(
                output_paths[kind]
            )

        txt_out.write_text(
            gerar_texto_base_v2(
                payload
            ),
            encoding="utf-8"
        )

        arquivos[
            "texto_base"
        ] = output_descriptor(
            txt_out
        )

    except Exception:
        # Remove somente arquivos reservados/criados
        # nesta execução.
        # Nunca remove arquivos anteriores.
        for path in reserved_paths:
            path.unlink(
                missing_ok=True
            )

        raise

    produtos_gerados = [
        (
            "CAPA"
            if kind == "capa"
            else "ESTRATEGICA"
        )
        for kind in selected_kinds
    ]

    return JSONResponse({
        "status": "ok",
        "versao": "v2.6",
        "codigo_ativo": payload.codigo_ativo,
        "tipo_saida": payload.tipo_saida,
        "produtos_gerados": produtos_gerados,
        "versao_saida": payload.versao_saida,
        "status_arquivo": payload.status_arquivo,
        "publicacao_automatica": (
            PUBLICACAO_AUTOMATICA
        ),
        "sobrescrita_silenciosa": False,
        "mensagem": (
            "Arquivos V2.6 gerados a partir "
            "dos modelos oficiais."
        ),
        "media_policy": payload.media_policy,
        "arquivos": arquivos,
        "resultados": resultados,
        "warnings": warnings,
    })
