from fastapi import FastAPI, HTTPException
from fastapi.responses import JSONResponse, FileResponse
from pydantic import BaseModel, Field
from typing import Optional, List, Dict, Any
from pathlib import Path
from pptx import Presentation
import shutil
import uuid
import unicodedata


app = FastAPI(
    title="Estate Gover Presentation Generator V2",
    version="2.0.0"
)

BASE_DIR = Path(__file__).resolve().parent
TEMPLATES_DIR = BASE_DIR / "templates"
OUTPUTS_DIR = BASE_DIR / "outputs"
OUTPUTS_DIR.mkdir(exist_ok=True)


class MediaItem(BaseModel):
    slot: str
    tipo: Optional[str] = None
    path: Optional[str] = None
    codigo_ativo: Optional[str] = None
    legenda: Optional[str] = None


class PresentationRequestV2(BaseModel):
    codigo_ativo: str
    nome_ativo: str
    localizacao: str

    tipo_ativo: Optional[str] = None
    area_aproximada: Optional[str] = None
    valor_referencia: Optional[str] = None

    descricao_capa: Optional[str] = None
    resumo_executivo: Optional[str] = None
    conheca_ativo: Optional[str] = None
    localizacao_texto: Optional[str] = None
    dimensao_status: Optional[str] = None
    potencial_urbanistico: Optional[str] = None
    diferenciais: Optional[str] = None
    escala_ativo: Optional[str] = None
    condicoes_negocio: Optional[str] = None
    observacoes_urbanisticas: Optional[str] = None

    # Campos legados da V1, mantidos por compatibilidade
    texto_base_capa: Optional[str] = None
    texto_base_estrategica: Optional[str] = None

    # V2
    midias: List[MediaItem] = Field(default_factory=list)
    media_policy: str = "asset_only"
    preservar_multimidia: bool = True
    formato_saida: str = "pptx"


PLACEHOLDER_FIELDS = [
    "codigo_ativo",
    "nome_ativo",
    "localizacao",
    "tipo_ativo",
    "area_aproximada",
    "valor_referencia",
    "descricao_capa",
    "resumo_executivo",
    "conheca_ativo",
    "localizacao_texto",
    "dimensao_status",
    "potencial_urbanistico",
    "diferenciais",
    "escala_ativo",
    "condicoes_negocio",
    "observacoes_urbanisticas",
    "texto_base_capa",
    "texto_base_estrategica",
]

FIXED_PREFIXES = (
    "EG_FIXED_",
    "EG_FIXED_LOGO",
    "EG_FIXED_QR",
    "EG_FIXED_WHATSAPP",
    "EG_FIXED_CONTATO",
)

ASSET_SLOT_PREFIXES = (
    "EG_SLOT_",
    "EG_ASSET_",
)


def normalize_filename(value: str) -> str:
    normalized = unicodedata.normalize("NFD", value)
    without_accents = normalized.encode("ascii", "ignore").decode("utf-8")
    return without_accents.lower()


def find_template(kind: str) -> Path:
    """
    Localiza o modelo oficial na pasta templates.

    A V2 aceita nomes oficiais e também os nomes atuais do repositório,
    sem quebrar a V1.
    """
    if kind == "capa":
        candidates = [
            "estate-gover-capa-modelo-final.pptx",
            "estate-governador-capa-modelo-final.pptx",
        ]
        keywords = ["capa"]
    elif kind == "estrategica":
        candidates = [
            "estate-gover-estrategica-modelo-final.pptx",
            "estate-gover-estratégica-modelo-final.pptx",
            "estate-governador-estrategica-modelo-final.pptx",
            "estate-governador-estratégica-modelo-final.pptx",
        ]
        keywords = ["estrategica", "estrateg"]
    else:
        raise ValueError("Tipo de modelo inválido.")

    for filename in candidates:
        path = TEMPLATES_DIR / filename
        if path.exists():
            return path

    for path in TEMPLATES_DIR.glob("*.pptx"):
        normalized = normalize_filename(path.name)
        if all(keyword in normalized for keyword in keywords):
            return path

    raise HTTPException(
        status_code=500,
        detail=f"Modelo oficial {kind} não encontrado na pasta templates."
    )


def safe_value(value: Any, fallback: str = "[pendente de confirmação]") -> str:
    if value is None:
        return fallback

    text = str(value).strip()
    if not text:
        return fallback

    return text


def build_placeholder_mapping(payload: PresentationRequestV2) -> Dict[str, str]:
    data = payload.dict()
    mapping = {}

    for field in PLACEHOLDER_FIELDS:
        mapping[f"{{{{{field}}}}}"] = safe_value(data.get(field))

    return mapping


def replace_placeholders_in_presentation(prs: Presentation, mapping: Dict[str, str]) -> int:
    """
    Substitui placeholders em caixas de texto existentes.

    Não altera layout, master, posição, tamanho, logos, rodapés, QR codes,
    botões ou objetos interativos.
    """
    replacements = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue

            if shape.text_frame is None:
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    original = run.text

                    for placeholder, value in mapping.items():
                        if placeholder in run.text:
                            run.text = run.text.replace(placeholder, value)

                    if run.text != original:
                        replacements += 1

    return replacements


def is_fixed_shape(shape) -> bool:
    name = getattr(shape, "name", "") or ""
    upper_name = name.upper()
    return any(upper_name.startswith(prefix.upper()) for prefix in FIXED_PREFIXES)


def get_asset_slot_name(shape) -> Optional[str]:
    name = getattr(shape, "name", "") or ""

    for prefix in ASSET_SLOT_PREFIXES:
        if name.upper().startswith(prefix.upper()):
            return name[len(prefix):]

    return None


def remove_shape(shape) -> None:
    element = shape._element
    element.getparent().remove(element)


def clear_text_shape(shape, message: str = "") -> bool:
    if not hasattr(shape, "text_frame") or shape.text_frame is None:
        return False

    shape.text_frame.clear()
    if message:
        shape.text = message

    return True


def apply_media_policy(
    prs: Presentation,
    payload: PresentationRequestV2,
    warnings: List[str]
) -> Dict[str, int]:
    """
    Política de mídia V2.

    Regra:
    - EG_FIXED_* deve ser preservado.
    - EG_SLOT_* ou EG_ASSET_* representa mídia do ativo.
    - Se houver mídia válida para o slot, ela entra no espaço existente.
    - Se não houver mídia válida, a mídia do ativo anterior é neutralizada.
    - Shapes sem marcação técnica não são removidos automaticamente, para evitar
      quebrar QR, WhatsApp, hyperlinks ou multimídia fixa.
    """
    stats = {
        "fixed_preserved": 0,
        "asset_slots_found": 0,
        "asset_slots_replaced": 0,
        "asset_slots_neutralized": 0,
        "untagged_media_untouched": 0,
    }

    midias_por_slot = {}

    for item in payload.midias:
        if item.codigo_ativo and item.codigo_ativo != payload.codigo_ativo:
            warnings.append(
                f"Mídia ignorada no slot {item.slot}: código do ativo diferente."
            )
            continue

        midias_por_slot[item.slot] = item

    for slide in prs.slides:
        for shape in list(slide.shapes):
            if is_fixed_shape(shape):
                stats["fixed_preserved"] += 1
                continue

            slot = get_asset_slot_name(shape)

            if not slot:
                continue

            stats["asset_slots_found"] += 1
            media = midias_por_slot.get(slot)

            if media and media.path:
                media_path = Path(media.path)

                if not media_path.exists():
                    warnings.append(
                        f"Mídia do slot {slot} não encontrada no caminho informado: {media.path}"
                    )
                    if clear_text_shape(shape, "[mídia do ativo não fornecida]"):
                        stats["asset_slots_neutralized"] += 1
                    continue

                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                remove_shape(shape)
                slide.shapes.add_picture(str(media_path), left, top, width=width, height=height)
                stats["asset_slots_replaced"] += 1
            else:
                # Neutralização segura para não vazar mídia de outro ativo.
                if clear_text_shape(shape, "[mídia do ativo não fornecida]"):
                    stats["asset_slots_neutralized"] += 1
                else:
                    remove_shape(shape)
                    stats["asset_slots_neutralized"] += 1

    if stats["asset_slots_found"] == 0:
        warnings.append(
            "Nenhum slot de mídia EG_SLOT_* ou EG_ASSET_* foi encontrado. "
            "A V2 não removeu mídia visual sem marcação para evitar quebrar QR, WhatsApp ou hyperlinks."
        )

    return stats


def gerar_texto_base(payload: PresentationRequestV2) -> str:
    return f"""ESTATE GOVER — TEXTO BASE V2

Código: {payload.codigo_ativo}
Ativo: {payload.nome_ativo}
Localização: {payload.localizacao}
Tipo: {payload.tipo_ativo or "[pendente de confirmação]"}
Área aproximada: {payload.area_aproximada or "[pendente de confirmação]"}
Valor de referência: {payload.valor_referencia or "[pendente de confirmação]"}

=== CAPA ===
{payload.descricao_capa or payload.texto_base_capa or "[pendente de confirmação]"}

=== ESTRATÉGICA ===

Resumo executivo:
{payload.resumo_executivo or payload.texto_base_estrategica or "[pendente de confirmação]"}

Conheça o ativo:
{payload.conheca_ativo or "[pendente de confirmação]"}

Localização:
{payload.localizacao_texto or "[pendente de confirmação]"}

Dimensão e status:
{payload.dimensao_status or "[pendente de confirmação]"}

Potencial construtivo e urbanístico:
{payload.potencial_urbanistico or "[necessário DM / consulta urbanística]"}

Diferenciais:
{payload.diferenciais or "[pendente de confirmação]"}

Escala do ativo:
{payload.escala_ativo or "[pendente de confirmação]"}

Condições de negócio:
{payload.condicoes_negocio or "[pendente de confirmação]"}

Observações urbanísticas:
{payload.observacoes_urbanisticas or "[necessário DM / consulta urbanística]"}
"""


def process_pptx(
    template_path: Path,
    output_path: Path,
    payload: PresentationRequestV2,
    warnings: List[str]
) -> Dict[str, Any]:
    shutil.copyfile(template_path, output_path)

    prs = Presentation(str(output_path))
    mapping = build_placeholder_mapping(payload)

    replacements = replace_placeholders_in_presentation(prs, mapping)
    media_stats = apply_media_policy(prs, payload, warnings)

    prs.save(str(output_path))

    if replacements == 0:
        warnings.append(
            f"Nenhum placeholder textual foi substituído em {output_path.name}. "
            "Confirme se o PPTX possui placeholders como {{codigo_ativo}}, {{localizacao}} e {{valor_referencia}}."
        )

    return {
        "template": template_path.name,
        "output": output_path.name,
        "text_replacements": replacements,
        "media_stats": media_stats,
    }


@app.get("/")
def healthcheck():
    return {
        "status": "ok",
        "service": "Estate Gover Presentation Generator",
        "version": "2.0.0"
    }


@app.get("/outputs/{filename}")
def baixar_arquivo(filename: str):
    file_path = OUTPUTS_DIR / filename

    if not file_path.exists():
        raise HTTPException(status_code=404, detail="Arquivo não encontrado.")

    return FileResponse(file_path)


@app.post("/v2/gerar-apresentacao-estate-gover")
def gerar_apresentacao_estate_gover_v2(payload: PresentationRequestV2):
    if payload.formato_saida != "pptx":
        raise HTTPException(
            status_code=400,
            detail="Apenas formato pptx é permitido."
        )

    warnings = []

    capa_template = find_template("capa")
    estrategica_template = find_template("estrategica")

    job_id = f"{payload.codigo_ativo}_{uuid.uuid4().hex[:8]}"

    capa_out = OUTPUTS_DIR / f"{job_id}_CAPA_V2.pptx"
    estrategica_out = OUTPUTS_DIR / f"{job_id}_ESTRATEGICA_V2.pptx"
    txt_out = OUTPUTS_DIR / f"{job_id}_texto_base_v2.txt"

    # Ordem obrigatória: CAPA antes da ESTRATÉGICA
    capa_result = process_pptx(capa_template, capa_out, payload, warnings)
    estrategica_result = process_pptx(estrategica_template, estrategica_out, payload, warnings)

    txt_out.write_text(gerar_texto_base(payload), encoding="utf-8")

    return JSONResponse({
        "status": "ok",
        "versao": "v2",
        "mensagem": "Arquivos V2 gerados a partir dos modelos oficiais.",
        "codigo_ativo": payload.codigo_ativo,
        "media_policy": payload.media_policy,
        "arquivos": {
            "capa_pptx_url": f"/outputs/{capa_out.name}",
            "estrategica_pptx_url": f"/outputs/{estrategica_out.name}",
            "texto_base_url": f"/outputs/{txt_out.name}"
        },
        "resultados": {
            "capa": capa_result,
            "estrategica": estrategica_result
        },
        "warnings": warnings
    })
