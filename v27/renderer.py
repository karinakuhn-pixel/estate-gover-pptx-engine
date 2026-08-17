"""Editable CAPA renderer for governed V2.7 homologation."""

from __future__ import annotations

import re
import zipfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt

from .models import PresentationRequestV27


CONCEPTUAL_DISCLAIMER = (
    "ILUSTRAÇÃO CONCEITUAL — NÃO REPRESENTA PROJETO APROVADO "
    "OU IMPLANTAÇÃO DEFINITIVA."
)


class RenderError(RuntimeError):
    pass


def governed_filename(payload: PresentationRequestV27) -> str:
    def part(value: str) -> str:
        text = value.strip().upper()
        text = re.sub(r"\s+", "-", text)
        text = re.sub(r"[^A-Z0-9À-Ü,.-]+", "-", text)
        return re.sub(r"-+", "-", text).strip("-_.")
    return (
        f"{part(payload.codigo_ativo)}_{part(payload.nome_ativo)}_"
        f"CAPA_{payload.versao_saida}_RASCUNHO.pptx"
    )


class CapaRendererV27:
    """Creates an editable PPTX without using legacy GBD/PPT templates."""

    def __init__(self, logo_path: str | Path | None = None):
        self.logo_path = Path(logo_path).resolve() if logo_path else None

    def _add_logo(self, slide) -> None:
        if self.logo_path and self.logo_path.is_file():
            slide.shapes.add_picture(
                str(self.logo_path),
                Inches(10.7),
                Inches(0.35),
                width=Inches(2.1),
            )

    @staticmethod
    def _add_title(slide, title: str, top=0.8, size=30) -> None:
        box = slide.shapes.add_textbox(
            Inches(0.75), Inches(top), Inches(11.2), Inches(0.8)
        )
        run = box.text_frame.paragraphs[0].add_run()
        run.text = title
        run.font.size = Pt(size)
        run.font.bold = True

    @staticmethod
    def _add_body(slide, text: str, top=1.8, height=4.6, size=18) -> None:
        box = slide.shapes.add_textbox(
            Inches(0.8), Inches(top), Inches(11.3), Inches(height)
        )
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)

    def render(self, payload: PresentationRequestV27, output_path: Path) -> dict[str, Any]:
        if payload.tipo_saida.value != "CAPA":
            raise RenderError("Fase 2 habilita somente CAPA")

        if any(
            image.classificacao == "CONCEITUAL"
            and image.autorizada
            and not image.ressalva
            for image in payload.imagens_autorizadas
        ):
            raise RenderError(
                "imagem CONCEITUAL autorizada exige ressalva explícita"
            )

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        # Abertura
        slide = prs.slides.add_slide(blank)
        self._add_logo(slide)
        self._add_title(slide, payload.nome_ativo, top=1.35, size=34)
        subtitle = f"{payload.municipio}/{payload.uf}"
        if payload.tese_criativa:
            subtitle += f"\n\n{payload.tese_criativa}"
        self._add_body(slide, subtitle, top=2.35, height=3.2, size=20)

        # Conteúdo governado
        for spec in payload.slides:
            slide = prs.slides.add_slide(blank)
            self._add_logo(slide)
            self._add_title(
                slide,
                spec.titulo or spec.tipo,
                top=0.85,
                size=28,
            )
            parts = [
                p for p in [
                    spec.mensagem_principal,
                    spec.texto_apoio,
                ] if p
            ]
            if spec.dados:
                parts.append(
                    "\n".join(f"{k}: {v}" for k, v in spec.dados.items())
                )
            if spec.ressalva:
                parts.append(f"RESSALVA: {spec.ressalva}")
            self._add_body(slide, "\n\n".join(parts) or " ", top=1.8)

        # Negociação
        neg = payload.negociacao
        if neg.preco_pedido or neg.condicao_comercial:
            slide = prs.slides.add_slide(blank)
            self._add_logo(slide)
            self._add_title(slide, "Negociação", top=0.85)
            lines = []
            if neg.preco_pedido:
                lines.append(f"Preço pedido: {neg.preco_pedido}")
            if neg.condicao_comercial:
                lines.append(f"Condições: {neg.condicao_comercial}")
            if neg.avaliacao:
                lines.append(f"Avaliação: {neg.avaliacao}")
            if neg.vgv:
                lines.append(f"VGV: {neg.vgv}")
            if neg.valor_residual:
                lines.append(f"Valor residual: {neg.valor_residual}")
            self._add_body(slide, "\n\n".join(lines), top=1.8)

        # Mapas e links
        if payload.mapas_links:
            slide = prs.slides.add_slide(blank)
            self._add_logo(slide)
            self._add_title(slide, "Localização e experiência espacial", top=0.85)
            y = 1.8
            for item in payload.mapas_links:
                box = slide.shapes.add_textbox(
                    Inches(0.8), Inches(y), Inches(11.2), Inches(0.6)
                )
                p = box.text_frame.paragraphs[0]
                r = p.add_run()
                r.text = f"{item.tipo} — {item.classificacao}"
                r.font.size = Pt(17)
                r.hyperlink.address = item.url
                y += 0.75

        # Contatos selecionados
        selected = [c for c in payload.contatos if c.selecionado]
        if selected:
            slide = prs.slides.add_slide(blank)
            self._add_logo(slide)
            self._add_title(slide, "Contatos", top=0.85)
            lines = []
            for contact in selected:
                row = [contact.nome]
                if contact.papel:
                    row.append(contact.papel)
                if contact.email:
                    row.append(contact.email)
                if contact.telefone:
                    row.append(contact.telefone)
                lines.append(" | ".join(row))
            self._add_body(slide, "\n\n".join(lines), top=1.8)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)
        qa = self.qa(output_path, payload)
        return {"output_path": str(output_path), **qa}

    def qa(self, path: Path, payload: PresentationRequestV27) -> dict[str, Any]:
        if not path.is_file() or path.stat().st_size <= 0:
            raise RenderError("PPTX não foi criado")
        if not zipfile.is_zipfile(path):
            raise RenderError("PPTX inválido")
        check = Presentation(path)
        if len(check.slides) < 1:
            raise RenderError("PPTX sem slides")

        hyperlinks_required = bool(payload.mapas_links)
        hyperlink_found = False
        for slide in check.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.hyperlink.address:
                            hyperlink_found = True

        if hyperlinks_required and not hyperlink_found:
            raise RenderError("hyperlink espacial obrigatório ausente")

        return {
            "slide_count": len(check.slides),
            "hyperlinks_ok": (not hyperlinks_required) or hyperlink_found,
            "editable": True,
        }
