import json
from pathlib import Path

import pytest
from fastapi import FastAPI
from fastapi.testclient import TestClient
from pptx import Presentation

from v27.audit import InMemoryAuditLogger, JsonlAuditLogger
from v27.middleware import V27OperationIdMiddleware
from v27.providers import LocalHomologationDriveProvider
from v27.renderer import CapaRendererV27, RenderError
from v27.routes import build_v27_router


FOLDER_ID = "1G9Ft-kriDZS4BYZBVmEsyWVq9XY5grxE"

PAYLOAD = {
    "codigo_ativo": "EG0029",
    "folder_id": FOLDER_ID,
    "nome_ativo": "AREA URBANA CARNIEL 5,2HA",
    "municipio": "Gramado",
    "uf": "RS",
    "tipo_saida": "CAPA",
    "versao_saida": "V01",
    "tese_criativa": "Uma área urbana em Gramado para discutir cenários — não para presumir um produto.",
    "slides": [
        {
            "tipo": "TESE",
            "titulo": "Território para discutir cenários",
            "mensagem_principal": "Leitura multivocacional, sem produto predeterminado.",
            "dados": {
                "matricula": "60.400",
                "area_registral": "52.201,88 m²",
                "referencia_comercial": "5,2 ha",
            },
        }
    ],
    "negociacao": {
        "preco_pedido": "R$ 65.000.000,00",
        "condicao_comercial": "Venda total ou entrada de R$ 5 milhões + permuta financeira a definir",
    },
    "mapas_links": [
        {
            "tipo": "Google Maps",
            "url": "https://maps.google.com/?q=-29.36086077,-50.85986071",
            "classificacao": "DELIMITAÇÃO COMERCIAL PRELIMINAR",
        }
    ],
    "contatos": [
        {
            "nome": "Karina Kuhn",
            "selecionado": True,
            "papel": "Contato técnico e institucional; bióloga e licenciamento ambiental",
            "email": "karina.kuhn@govertrust.com.br",
        }
    ],
    "formato_saida": "pptx",
}


def make_asset(tmp_path):
    root = tmp_path / "EG0029"
    root.mkdir(parents=True)
    (root / "04_CAPA").mkdir()
    (root / "05_ESTRATEGICA").mkdir()
    return root


def make_client(tmp_path, logo=None, audit=None):
    root = make_asset(tmp_path)
    provider = LocalHomologationDriveProvider({FOLDER_ID: root})
    audit = audit or InMemoryAuditLogger()
    renderer = CapaRendererV27(logo_path=logo)
    app = FastAPI()
    app.add_middleware(V27OperationIdMiddleware, audit_logger=audit)
    app.include_router(build_v27_router(audit, provider, renderer))
    return TestClient(app), root, audit


def test_phase2_generates_editable_capa_and_uploads_to_04_capa(tmp_path):
    client, root, audit = make_client(tmp_path)
    response = client.post("/v2.7/gerar-apresentacao-estate-gover", json=PAYLOAD)
    assert response.status_code == 200, response.text
    body = response.json()
    assert body["status"] == "rascunho_gerado"
    assert body["status_arquivo"] == "RASCUNHO"
    assert body["publicacao_automatica"] is False
    assert body["tipo_saida"] == "CAPA"
    assert len(body["arquivos"]) == 1

    pptx = root / "04_CAPA" / body["arquivos"][0]["arquivo"]
    assert pptx.exists()
    prs = Presentation(pptx)
    assert len(prs.slides) >= 4
    assert any(e["event"] == "upload_completed" for e in audit.events)


def test_conflict_returns_409_and_preserves_existing_file(tmp_path):
    client, root, _ = make_client(tmp_path)
    first = client.post("/v2.7/gerar-apresentacao-estate-gover", json=PAYLOAD)
    assert first.status_code == 200
    target = root / "04_CAPA" / first.json()["arquivos"][0]["arquivo"]
    original = target.read_bytes()

    second = client.post("/v2.7/gerar-apresentacao-estate-gover", json=PAYLOAD)
    assert second.status_code == 409
    assert target.read_bytes() == original


def test_conflict_is_checked_again_immediately_before_upload(tmp_path):
    root = make_asset(tmp_path)

    class AppearingConflictProvider(LocalHomologationDriveProvider):
        def __init__(self):
            super().__init__({FOLDER_ID: root})
            self.checks = 0
            self.upload_called = False

        def file_exists(self, folder_id, filename):
            self.checks += 1
            return self.checks >= 2

        def upload_draft(self, folder_id, path):
            self.upload_called = True
            return super().upload_draft(folder_id, path)

    provider = AppearingConflictProvider()
    audit = InMemoryAuditLogger()
    app = FastAPI()
    app.add_middleware(V27OperationIdMiddleware, audit_logger=audit)
    app.include_router(build_v27_router(audit, provider, CapaRendererV27()))

    response = TestClient(app).post(
        "/v2.7/gerar-apresentacao-estate-gover", json=PAYLOAD
    )
    assert response.status_code == 409
    assert provider.checks == 2
    assert provider.upload_called is False
    assert any(event["event"] == "pre_upload_conflict" for event in audit.events)


def test_unknown_folder_id_is_blocked(tmp_path):
    client, _, audit = make_client(tmp_path)
    response = client.post(
        "/v2.7/gerar-apresentacao-estate-gover",
        json={**PAYLOAD, "folder_id": "nao-autorizado"},
    )
    assert response.status_code == 422
    assert any(e["event"] == "phase2_failed" for e in audit.events)


def test_missing_04_capa_is_blocked_and_not_created(tmp_path):
    root = tmp_path / "EG0029"
    root.mkdir(parents=True)
    (root / "05_ESTRATEGICA").mkdir()
    provider = LocalHomologationDriveProvider({FOLDER_ID: root})
    audit = InMemoryAuditLogger()
    app = FastAPI()
    app.add_middleware(V27OperationIdMiddleware, audit_logger=audit)
    app.include_router(build_v27_router(audit, provider, CapaRendererV27()))
    client = TestClient(app)

    response = client.post("/v2.7/gerar-apresentacao-estate-gover", json=PAYLOAD)
    assert response.status_code == 422
    assert not (root / "04_CAPA").exists()


def test_estrategica_stays_blocked_in_phase2(tmp_path):
    client, _, _ = make_client(tmp_path)
    response = client.post(
        "/v2.7/gerar-apresentacao-estate-gover",
        json={**PAYLOAD, "tipo_saida": "ESTRATEGICA"},
    )
    assert response.status_code == 501


def test_map_hyperlink_is_present_in_pptx(tmp_path):
    client, root, _ = make_client(tmp_path)
    response = client.post("/v2.7/gerar-apresentacao-estate-gover", json=PAYLOAD)
    path = root / "04_CAPA" / response.json()["arquivos"][0]["arquivo"]
    prs = Presentation(path)
    links = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        if run.hyperlink.address:
                            links.append(run.hyperlink.address)
    assert PAYLOAD["mapas_links"][0]["url"] in links


def test_conceptual_image_requires_explicit_disclaimer(tmp_path):
    client, _, _ = make_client(tmp_path)
    response = client.post(
        "/v2.7/gerar-apresentacao-estate-gover",
        json={
            **PAYLOAD,
            "imagens_autorizadas": [{
                "referencia": "conceito.png",
                "classificacao": "CONCEITUAL",
                "autorizada": True,
            }],
        },
    )
    assert response.status_code == 422
    assert "ressalva" in response.json()["detail"]


def test_jsonl_audit_is_append_only_and_keeps_operation_id(tmp_path):
    audit_path = tmp_path / "audit" / "v27.jsonl"
    audit = JsonlAuditLogger(audit_path)
    client, _, _ = make_client(tmp_path / "asset", audit=audit)
    response = client.post("/v2.7/gerar-apresentacao-estate-gover", json=PAYLOAD)
    assert response.status_code == 200
    lines = [json.loads(line) for line in audit_path.read_text(encoding="utf-8").splitlines()]
    assert len(lines) >= 4
    ids = {line.get("operation_id") for line in lines if line.get("operation_id")}
    assert ids == {response.json()["operation_id"]}
    assert all(line["version"] == "v2.7" for line in lines)


def test_logo_file_can_be_embedded_without_becoming_template(tmp_path):
    # Usa uma imagem mínima válida criada para o teste; renderer não depende de template legado.
    from PIL import Image
    logo = tmp_path / "logo.png"
    Image.new("RGB", (100, 40), "white").save(logo)
    client, root, _ = make_client(tmp_path / "case", logo=logo)
    response = client.post("/v2.7/gerar-apresentacao-estate-gover", json=PAYLOAD)
    assert response.status_code == 200
    path = root / "04_CAPA" / response.json()["arquivos"][0]["arquivo"]
    prs = Presentation(path)
    assert any(slide.shapes for slide in prs.slides)
